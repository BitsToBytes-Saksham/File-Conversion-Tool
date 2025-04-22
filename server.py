import socket
import os
import tempfile
import shutil
import zipfile
import io
import comtypes.client # type: ignore
import pypdf # Use pypdf for newer features / potentially better handling
import traceback # For detailed error logging

# --- Pillow for image saving (still needed) ---
from PIL import Image # type: ignore

# --- PyMuPDF (fitz) for PDF rendering (replaces pdf2image) ---
import fitz # PyMuPDF

# --- Other libraries for specific conversions/actions ---
from pdf2docx import Converter as PDF2WordConverter # For PDF to Word
from pptx import Presentation # For PDF to PPTX
from pptx.util import Inches, Emu # For PDF to PPTX (Emu for direct use)
from reportlab.pdfgen import canvas # For Page Numbers
from reportlab.lib.units import inch # For Page Numbers
# from reportlab.lib.pagesizes import letter # Use actual page size instead
# --- End Imports ---


# --- Office Conversion Functions (Require MS Office Installed) ---
# Ensure COM is initialized/uninitialized per function call for thread safety
def convert_docx_to_pdf(input_path, output_path):
    word = None
    doc = None
    com_initialized = False
    try:
        try:
            comtypes.CoInitialize()
            com_initialized = True
        except OSError: # Already initialized on this thread
            com_initialized = False
        word = comtypes.client.CreateObject('Word.Application')
        word.Visible = False
        doc = word.Documents.Open(input_path)
        doc.SaveAs(output_path, FileFormat=17) # 17 = wdFormatPDF
        print(f"Successfully converted DOCX: {input_path} to {output_path}")
    except Exception as e:
        print(f"Error converting DOCX to PDF: {e}\n{traceback.format_exc()}")
        raise # Re-raise the exception
    finally:
        if doc: doc.Close(False)
        if word: word.Quit()
        word, doc = None, None # Release objects
        if com_initialized: comtypes.CoUninitialize()

def convert_pptx_to_pdf(input_path, output_path):
    powerpoint = None
    presentation = None
    com_initialized = False
    try:
        try:
            comtypes.CoInitialize()
            com_initialized = True
        except OSError:
            com_initialized = False
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        presentation.SaveAs(output_path, 32) # 32 = ppSaveAsPDF
        print(f"Successfully converted PPTX: {input_path} to {output_path}")
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}\n{traceback.format_exc()}")
        raise
    finally:
        if presentation: presentation.Close()
        if powerpoint: powerpoint.Quit()
        powerpoint, presentation = None, None
        if com_initialized: comtypes.CoUninitialize()

def convert_xlsx_to_pdf(input_path, output_path):
    excel = None
    wb = None
    com_initialized = False
    try:
        try:
            comtypes.CoInitialize()
            com_initialized = True
        except OSError:
            com_initialized = False
        excel = comtypes.client.CreateObject("Excel.Application")
        excel.Visible = False
        wb = excel.Workbooks.Open(input_path)
        wb.ExportAsFixedFormat(0, output_path) # 0 = xlTypePDF
        print(f"Successfully converted XLSX: {input_path} to {output_path}")
    except Exception as e:
        print(f"Error converting XLSX to PDF: {e}\n{traceback.format_exc()}")
        raise
    finally:
        if wb: wb.Close(False)
        if excel: excel.Quit()
        excel, wb = None, None
        if com_initialized: comtypes.CoUninitialize()

def convert_html_to_pdf(input_path, output_path):
    word = None
    doc = None
    com_initialized = False
    try:
        try:
            comtypes.CoInitialize()
            com_initialized = True
        except OSError:
            com_initialized = False
        word = comtypes.client.CreateObject('Word.Application')
        word.Visible = False
        doc = word.Documents.Open(input_path, Format="wdOpenFormatWebPages")
        doc.SaveAs(output_path, FileFormat=17) # wdFormatPDF
        print(f"Successfully converted HTML: {input_path} to {output_path}")
    except Exception as e:
        print(f"Error converting HTML to PDF: {e}\n{traceback.format_exc()}")
        raise
    finally:
        if doc: doc.Close(False)
        if word: word.Quit()
        word, doc = None, None
        if com_initialized: comtypes.CoUninitialize()

# --- General File Type Handler ---
def handle_file_conversion(ext, input_path, output_path):
    """Handles conversion FROM various formats TO PDF."""
    ext = ext.lower()
    abs_input_path = os.path.abspath(input_path)
    abs_output_path = os.path.abspath(output_path)
    print(f"Handling conversion for {ext}: {abs_input_path} -> {abs_output_path}")

    try:
        if ext in [".jpg", ".jpeg", ".png"]:
            with Image.open(abs_input_path) as img:
                img_to_save = img.convert('RGB') if img.mode in ['RGBA', 'P'] else img
                img_to_save.save(abs_output_path, "PDF", resolution=100.0, save_all=False) # Basic save
            print(f"Successfully converted Image: {abs_input_path} to {abs_output_path}")
        elif ext == ".docx":
            convert_docx_to_pdf(abs_input_path, abs_output_path)
        elif ext == ".pptx":
            convert_pptx_to_pdf(abs_input_path, abs_output_path)
        elif ext == ".xlsx":
            convert_xlsx_to_pdf(abs_input_path, abs_output_path)
        elif ext == ".html":
            convert_html_to_pdf(abs_input_path, abs_output_path)
        else:
            raise ValueError(f"Unsupported file format for 'convert' action: {ext}")
    except Exception as e:
         print(f"Failure during handle_file_conversion for {ext}: {e}")
         raise # Re-raise to be caught by main handler

# --- PDF Security Functions ---
def encrypt_pdf(input_path, output_path, password):
    """Encrypts a PDF using pypdf."""
    try:
        reader = pypdf.PdfReader(input_path)
        if reader.is_encrypted:
             raise ValueError("Cannot re-encrypt an already encrypted PDF. Decrypt first.")

        writer = pypdf.PdfWriter()
        writer.clone_document_from_reader(reader) # Clone pages and metadata
        writer.encrypt(password)

        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        print(f"Successfully encrypted: {input_path} to {output_path}")
        return True # Indicate success
    except Exception as e:
        print(f"Encryption error: {e}\n{traceback.format_exc()}")
        # Do not return False here, let the exception propagate
        raise ValueError(f"PDF Encryption failed: {e}") # Raise specific error


def decrypt_pdf(input_path, output_path, password):
    """Decrypts a PDF using pypdf."""
    try:
        reader = pypdf.PdfReader(input_path)
        if reader.is_encrypted:
            print("PDF is encrypted, attempting decryption...")
            decrypt_result = reader.decrypt(password)
            # Check the result type for success indication (pypdf >= 3.x)
            if decrypt_result == pypdf.PasswordType.OWNER_PASSWORD:
                 print("Decrypted using owner password.")
            elif decrypt_result == pypdf.PasswordType.USER_PASSWORD:
                 print("Decrypted using user password.")
            elif decrypt_result == pypdf.PasswordType.NOT_DECRYPTED:
                 # Raise a specific error for incorrect password
                 raise ValueError("Incorrect password provided for decryption.")
            else:
                 print(f"Decryption status: {decrypt_result}") # Log other statuses
        else:
            print("PDF was not encrypted, copying file as is.")

        writer = pypdf.PdfWriter()
        writer.clone_document_from_reader(reader) # Clone regardless of encryption status

        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        print(f"Successfully decrypted/copied: {input_path} to {output_path}")
        return True # Indicate success
    except ValueError as ve: # Catch specific password error
        print(f"Decryption failed: {ve}")
        raise # Re-raise the specific error
    except Exception as e:
        print(f"Decryption error: {e}\n{traceback.format_exc()}")
        raise ValueError(f"PDF Decryption failed: {e}") # Raise specific error

# --- PDF to Other Format Functions ---
def convert_pdf_to_jpg(input_path, output_dir):
    """Converts PDF pages to JPG images using PyMuPDF."""
    created_files = []
    doc = None
    try:
        os.makedirs(output_dir, exist_ok=True)
        print(f"Converting PDF to JPG using PyMuPDF: {input_path} -> {output_dir}")
        doc = fitz.open(input_path)
        if not doc.page_count: raise ValueError("Input PDF has no pages.")
        dpi = 150
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=dpi)
            output_jpg_path = os.path.join(output_dir, f"page_{i+1:03d}.jpg")
            pix.save(output_jpg_path, "jpeg")
            created_files.append(output_jpg_path)
        doc.close()
        doc = None
        if not created_files: raise Exception("PyMuPDF failed to convert any pages.")
        print(f"Converted {len(created_files)} pages to JPG in {output_dir}")
        return created_files
    except Exception as e:
        print(f"Error converting PDF to JPG with PyMuPDF: {e}\n{traceback.format_exc()}")
        if doc: doc.close() # Ensure closure on error
        if os.path.exists(output_dir): shutil.rmtree(output_dir, ignore_errors=True)
        raise # Re-raise
    finally:
         if doc: doc.close()

def convert_pdf_to_word(input_path, output_path):
    """Converts PDF to DOCX using pdf2docx."""
    try:
        print(f"Converting PDF to DOCX: {input_path} -> {output_path}")
        cv = PDF2WordConverter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        print("Successfully converted PDF to DOCX")
    except Exception as e:
        print(f"Error converting PDF to Word: {e}\n{traceback.format_exc()}")
        raise

def convert_pdf_to_pptx(input_path, output_path, temp_img_dir):
    """Converts PDF pages to images (PyMuPDF) and inserts into PPTX."""
    doc = None
    try:
        print(f"Converting PDF to PPTX (as images using PyMuPDF): {input_path} -> {output_path}")
        os.makedirs(temp_img_dir, exist_ok=True)
        doc = fitz.open(input_path)
        if not doc.page_count: raise ValueError("Input PDF has no pages.")

        image_paths = []
        dpi = 150
        for i, page in enumerate(doc):
             pix = page.get_pixmap(dpi=dpi)
             img_path = os.path.join(temp_img_dir, f"slide_{i+1:03d}.png")
             pix.save(img_path, "png")
             image_paths.append(img_path)
        doc.close(); doc = None
        if not image_paths: raise Exception("PyMuPDF failed to convert pages for PPTX.")

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[5]
        slide_width_emu, slide_height_emu = prs.slide_width, prs.slide_height

        for i, img_path in enumerate(image_paths):
            if not os.path.exists(img_path): continue
            slide = prs.slides.add_slide(blank_slide_layout)
            try:
                with Image.open(img_path) as im: img_width_px, img_height_px = im.size
                img_width_emu = Emu(img_width_px * 914400 / 96)
                img_height_emu = Emu(img_height_px * 914400 / 96)
                ratio = min(slide_width_emu / img_width_emu, slide_height_emu / img_height_emu) if img_width_emu > 0 and img_height_emu > 0 else 1
                pic_width_emu = int(img_width_emu * ratio)
                pic_height_emu = int(img_height_emu * ratio)
                left = int((slide_width_emu - pic_width_emu) / 2)
                top = int((slide_height_emu - pic_height_emu) / 2)
                slide.shapes.add_picture(img_path, left, top, width=pic_width_emu, height=pic_height_emu)
            except Exception as pic_e: print(f"Error adding picture {img_path} to slide {i+1}: {pic_e}")

        prs.save(output_path)
        print(f"Successfully created PPTX with {len(image_paths)} slides.")
    except Exception as e:
        print(f"Error converting PDF to PPTX: {e}\n{traceback.format_exc()}")
        if doc: doc.close()
        raise
    finally:
        if os.path.exists(temp_img_dir): shutil.rmtree(temp_img_dir, ignore_errors=True)


# --- PDF Manipulation Functions ---
def compress_pdf(input_path, output_path):
    """Compresses PDF streams using pypdf."""
    try:
        print(f"Compressing PDF: {input_path} -> {output_path}")
        reader = pypdf.PdfReader(input_path)
        writer = pypdf.PdfWriter()
        writer.clone_document_from_reader(reader)
        for page in writer.pages:
             try: # Compression can sometimes fail on complex/corrupt pages
                  page.compress_content_streams()
             except Exception as comp_e:
                  page_index = writer.get_page_number(page) # Get index for logging
                  print(f"Warning: Could not compress content stream for page {page_index + 1}: {comp_e}")
        with open(output_path, 'wb') as f: writer.write(f)
        print("Finished compress operation.")
    except Exception as e:
        print(f"Error compressing PDF: {e}\n{traceback.format_exc()}")
        raise

def split_pdf(input_path, output_dir, ranges_str):
    """Splits PDF based on page ranges using pypdf."""
    output_files = []
    try:
        print(f"Splitting PDF: {input_path} based on ranges '{ranges_str}'")
        os.makedirs(output_dir, exist_ok=True)
        reader = pypdf.PdfReader(input_path)
        num_pages = len(reader.pages)
        if num_pages == 0: raise ValueError("Input PDF has no pages to split.")

        parts = ranges_str.split(',')
        file_index = 1
        for part in parts:
            part = part.strip()
            if not part: continue
            writer = pypdf.PdfWriter() # Create new writer for each part
            pages_added = False
            try: # Add try block for range parsing/page adding
                if '-' in part:
                    start_str, end_str = part.split('-', 1)
                    start = int(start_str) if start_str else 1
                    end = int(end_str) if end_str else num_pages
                    if not (1 <= start <= end <= num_pages): raise ValueError(f"Range {start}-{end} out of bounds (1-{num_pages})")
                    for i in range(start - 1, end): writer.add_page(reader.pages[i]); pages_added = True
                else:
                    page_num = int(part)
                    if not (1 <= page_num <= num_pages): raise ValueError(f"Page number {page_num} out of bounds (1-{num_pages})")
                    writer.add_page(reader.pages[page_num - 1]); pages_added = True

                if pages_added:
                    safe_base = "".join(c if c.isalnum() else "_" for c in os.path.basename(input_path))
                    output_filename = os.path.join(output_dir, f"split_{file_index}_{safe_base}.pdf")
                    with open(output_filename, 'wb') as f: writer.write(f)
                    output_files.append(output_filename)
                    print(f"Created split file: {output_filename}")
                    file_index += 1
                else: print(f"Warning: Range '{part}' resulted in no pages being added.") # Should not happen if logic is correct
            except (ValueError, IndexError) as parse_err:
                print(f"Warning: Skipping invalid range/page '{part}': {parse_err}")
            except Exception as page_err:
                 print(f"Warning: Error processing range '{part}': {page_err}")

        if not output_files: raise ValueError("No valid output files created. Check input ranges and PDF.")
        return output_files
    except Exception as e:
        print(f"Error splitting PDF: {e}\n{traceback.format_exc()}")
        # Cleanup might be complex if some files were created
        raise


def merge_pdfs(input_zip_path, output_path):
    """Merges multiple PDFs from a zip file using pypdf."""
    temp_extract_dir = os.path.join(tempfile.gettempdir(), f"merge_{os.path.basename(output_path)}_extracted_{os.getpid()}")
    merger = None
    extracted_files = [] # Keep track for cleanup
    try:
        print(f"Merging PDFs from zip: {input_zip_path} -> {output_path}")
        os.makedirs(temp_extract_dir, exist_ok=True)

        pdf_files_to_merge = []
        with zipfile.ZipFile(input_zip_path, 'r') as zip_ref:
            for member in zip_ref.infolist():
                if member.is_dir() or not member.filename.lower().endswith('.pdf') or '../' in member.filename: continue
                target_path = os.path.join(temp_extract_dir, os.path.basename(member.filename))
                try:
                     with zip_ref.open(member.filename) as source, open(target_path, "wb") as target:
                         shutil.copyfileobj(source, target)
                     extracted_files.append(target_path) # Add to cleanup list
                     pdf_files_to_merge.append(target_path)
                     print(f"Extracted for merging: {os.path.basename(target_path)}")
                except Exception as extract_e: print(f"Warning: Failed to extract {member.filename}: {extract_e}")

        if not pdf_files_to_merge: raise ValueError("No valid PDF files found or extracted from the zip archive.")
        pdf_files_to_merge.sort() # Sort for predictable order

        merger = pypdf.PdfMerger()
        print("Appending PDFs...")
        for pdf_path in pdf_files_to_merge:
            try: merger.append(pdf_path); print(f"Appended: {os.path.basename(pdf_path)}")
            except Exception as append_e: print(f"Warning: Could not append PDF '{os.path.basename(pdf_path)}'. Skipping. Error: {append_e}")

        if len(merger.pages) == 0: raise Exception("Merging resulted in an empty PDF. Check input files and logs.")

        print("Writing merged PDF...")
        with open(output_path, 'wb') as f: merger.write(f)
        merger.close(); merger = None # Close explicitly
        print("Successfully merged PDFs.")
    except Exception as e:
        print(f"Error merging PDFs: {e}\n{traceback.format_exc()}")
        if merger: merger.close() # Ensure close on error
        raise
    finally:
        if os.path.exists(temp_extract_dir):
            print(f"Cleaning up temporary merge directory: {temp_extract_dir}")
            shutil.rmtree(temp_extract_dir, ignore_errors=True)


def rotate_pdf(input_path, output_path, pages_str, angle):
    """Rotates specified pages in a PDF using pypdf."""
    try:
        print(f"Rotating PDF: {input_path}, pages '{pages_str}', angle {angle}")
        angle = int(angle)
        if angle % 90 != 0: raise ValueError("Rotation angle must be a multiple of 90.")

        reader = pypdf.PdfReader(input_path)
        writer = pypdf.PdfWriter()
        num_pages = len(reader.pages)
        if num_pages == 0: raise ValueError("Input PDF has no pages to rotate.")

        rotate_indices = set()
        pages_str_lower = pages_str.strip().lower()
        if pages_str_lower == 'all':
            rotate_indices = set(range(num_pages))
        elif pages_str_lower:
            parts = pages_str_lower.split(',')
            for part in parts:
                part = part.strip();
                if not part: continue
                try:
                    if '-' in part:
                        start, end = part.split('-', 1)
                        start = int(start) if start else 1
                        end = int(end) if end else num_pages
                        if not (1 <= start <= end <= num_pages): raise ValueError(f"Range {start}-{end} out of bounds (1-{num_pages})")
                        rotate_indices.update(range(start - 1, end))
                    else:
                        page_num = int(part)
                        if not (1 <= page_num <= num_pages): raise ValueError(f"Page number {page_num} out of bounds (1-{num_pages})")
                        rotate_indices.add(page_num - 1)
                except ValueError as parse_err: raise ValueError(f"Invalid page specification '{part}': {parse_err}")

        if not rotate_indices: print("Warning: No pages selected for rotation based on input.")

        writer.clone_document_from_reader(reader) # Clone first
        rotated_count = 0
        for i in rotate_indices:
             if 0 <= i < len(writer.pages):
                 writer.pages[i].rotate(angle)
                 rotated_count += 1
             else: print(f"Warning: Index {i} from selection is out of bounds, skipping.")
        print(f"Rotated {rotated_count} pages by {angle} degrees")

        with open(output_path, 'wb') as f: writer.write(f)
        print("Successfully completed rotate operation.")
    except Exception as e:
        print(f"Error rotating PDF: {e}\n{traceback.format_exc()}")
        raise


def add_page_numbers_to_pdf(input_path, output_path, position):
    """Adds page numbers to a PDF using ReportLab and pypdf."""
    try:
        print(f"Adding page numbers ({position}) to: {input_path} -> {output_path}")
        reader = pypdf.PdfReader(input_path)
        writer = pypdf.PdfWriter()
        num_pages = len(reader.pages)
        if num_pages == 0: raise ValueError("Input PDF has no pages to add numbers to.")

        added_count = 0
        for i, page in enumerate(reader.pages):
            packet = io.BytesIO()
            try:
                page_width = float(page.mediabox.width)
                page_height = float(page.mediabox.height)
                can = canvas.Canvas(packet, pagesize=(page_width, page_height))

                page_num_text = f"Page {i + 1} of {num_pages}"
                font_size = 9
                can.setFont("Helvetica", font_size)
                text_width = can.stringWidth(page_num_text, "Helvetica", font_size)
                margin = 0.5 * inch

                x = (page_width - text_width) / 2; y = margin # Default bottom-center
                if position == 'bottom-left': x = margin
                elif position == 'bottom-right': x = page_width - text_width - margin
                elif position == 'top-center': y = page_height - margin - font_size
                elif position == 'top-left': x = margin; y = page_height - margin - font_size
                elif position == 'top-right': x = page_width - text_width - margin; y = page_height - margin - font_size

                can.drawString(x, y, page_num_text)
                can.save()
                packet.seek(0)

                overlay_pdf = pypdf.PdfReader(packet)
                if overlay_pdf.pages:
                     page.merge_page(overlay_pdf.pages[0])
                     added_count += 1
                else: print(f"Warning: ReportLab overlay for page {i+1} was empty.")
            except Exception as overlay_e:
                 print(f"Warning: Could not create or merge overlay for page {i+1}: {overlay_e}")
                 # Add the original page anyway
            finally:
                 writer.add_page(page) # Add original or merged page

        print(f"Added page numbers to {added_count} pages.")
        with open(output_path, "wb") as f: writer.write(f)
        print("Successfully completed add page numbers operation.")
    except Exception as e:
        print(f"Error adding page numbers: {e}\n{traceback.format_exc()}")
        raise


# --- Main Server Logic ---
def server_program():
    HOST = '127.0.0.1'
    PORT = 65432

    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1) # Allow address reuse
        s.bind((HOST, PORT))
        s.listen()
        print(f"Server listening on {HOST}:{PORT}")

        while True:
            conn, addr = s.accept()
            # Set timeouts for the connection
            conn.settimeout(60.0) # General timeout
            with conn:
                print(f"\n--- New connection from {addr} ---")
                input_path = None
                output_path = None
                temp_dir_for_action = None
                received_zip_path = None
                action_success = False
                action = "unknown" # Default action for logging errors early

                try:
                    # 1. Get Action
                    action_bytes = conn.recv(1024)
                    if not action_bytes: raise ConnectionAbortedError("Client disconnected before sending action.")
                    action = action_bytes.decode().strip()
                    conn.sendall(b"ACK_ACTION")
                    print(f"Received action: {action}")

                    # 2. Get Input Filename
                    base_filename_bytes = conn.recv(1024)
                    if not base_filename_bytes: raise ConnectionAbortedError("Client disconnected before sending filename.")
                    base_filename = base_filename_bytes.decode()
                    conn.sendall(b"ACK_FILENAME")
                    print(f"Received base filename: {base_filename}")

                    # --- Setup Paths ---
                    temp_dir = tempfile.gettempdir()
                    input_filename_base, input_ext = os.path.splitext(base_filename)
                    safe_base = "".join(c if c.isalnum() or c in ['_', '-'] else '_' for c in input_filename_base)[:50] # Limit length
                    pid_suffix = f"_{os.getpid()}"

                    if action == "merge":
                         input_path = os.path.join(temp_dir, f"{safe_base}{pid_suffix}.zip")
                         received_zip_path = input_path
                    else:
                         input_path = os.path.join(temp_dir, f"{safe_base}{pid_suffix}{input_ext}")

                    # Define output path/suggestion (will be created by the action function)
                    output_filename_suggestion = f"{safe_base}_processed.pdf" # Default
                    if action == "pdf_to_jpg":
                        temp_dir_for_action = os.path.join(temp_dir, f"{safe_base}_jpg_output{pid_suffix}")
                        output_path = os.path.join(temp_dir, f"{safe_base}_images{pid_suffix}.zip")
                        output_filename_suggestion = f"{safe_base}_images.zip"
                    # ... (other specific actions as before) ...
                    elif action == "pdf_to_word":
                        output_path = os.path.join(temp_dir, f"{safe_base}{pid_suffix}.docx")
                        output_filename_suggestion = f"{safe_base}.docx"
                    elif action == "pdf_to_pptx":
                         temp_dir_for_action = os.path.join(temp_dir, f"{safe_base}_pptx_temp_imgs{pid_suffix}")
                         output_path = os.path.join(temp_dir, f"{safe_base}{pid_suffix}.pptx")
                         output_filename_suggestion = f"{safe_base}.pptx"
                    elif action == "split":
                        temp_dir_for_action = os.path.join(temp_dir, f"{safe_base}_split_output{pid_suffix}")
                        output_path = os.path.join(temp_dir, f"{safe_base}_split_files{pid_suffix}.zip")
                        output_filename_suggestion = f"{safe_base}_split_files.zip"
                    elif action == "merge":
                         output_filename_base = f"merged_{input_filename_base}"[:50] # Limit merged name length
                         output_path = os.path.join(temp_dir, f"{output_filename_base}{pid_suffix}.pdf")
                         output_filename_suggestion = f"{output_filename_base}.pdf"
                    else: # Standard PDF output actions
                        suffix = action if action != "convert" else "processed"
                        output_filename_base = f"{safe_base}_{suffix}"
                        output_path = os.path.join(temp_dir, f"{output_filename_base}{pid_suffix}.pdf")
                        output_filename_suggestion = f"{output_filename_base}.pdf"
                    # --- End Path Setup ---

                    # 3. Get File Size
                    size_bytes = conn.recv(16)
                    if not size_bytes: raise ConnectionAbortedError("Client disconnected before sending size.")
                    size = int(size_bytes.decode().strip())
                    conn.sendall(b"ACK_SIZE")
                    print(f"Expecting file size: {size} bytes")

                    # 4. Receive File Data
                    print(f"Receiving data into: {input_path}...")
                    received_bytes = 0
                    with open(input_path, "wb") as f:
                        while received_bytes < size:
                            # Adjust timeout for potentially large file transfer
                            conn.settimeout(120.0) # 2 minutes per chunk read?
                            chunk = conn.recv(min(65536, size - received_bytes))
                            if not chunk: raise ConnectionAbortedError(f"Client disconnected during file transfer ({received_bytes}/{size} received).")
                            f.write(chunk)
                            received_bytes += len(chunk)
                    conn.settimeout(600.0) # Reset longer timeout for processing
                    print(f"Received {received_bytes} bytes and saved to {input_path}")
                    if received_bytes != size: print(f"Warning: Received bytes ({received_bytes}) differs from expected size ({size}).")

                    # --- Receive Extra Options ---
                    options = {} # Dictionary to hold options
                    if action in ("encrypt", "decrypt"):
                        print("Waiting for password...")
                        pass_bytes = conn.recv(1024); conn.sendall(b"ACK_PASS")
                        if not pass_bytes: raise ConnectionAbortedError("No password received.")
                        options["password"] = pass_bytes.decode()
                        print("Received password.")
                    elif action == "split":
                        print("Waiting for split ranges...")
                        ranges_bytes = conn.recv(1024); conn.sendall(b"ACK_RANGES")
                        if not ranges_bytes: raise ConnectionAbortedError("No split ranges received.")
                        options["ranges"] = ranges_bytes.decode()
                        print(f"Received ranges: {options['ranges']}")
                    elif action == "rotate":
                        print("Waiting for rotate pages...")
                        pages_bytes = conn.recv(1024); conn.sendall(b"ACK_PAGES")
                        if not pages_bytes: raise ConnectionAbortedError("No rotate pages received.")
                        options["pages"] = pages_bytes.decode()
                        print(f"Received pages: {options['pages']}")
                        print("Waiting for rotate angle...")
                        angle_bytes = conn.recv(1024); conn.sendall(b"ACK_ANGLE")
                        if not angle_bytes: raise ConnectionAbortedError("No rotate angle received.")
                        options["angle"] = angle_bytes.decode()
                        print(f"Received angle: {options['angle']}")
                    elif action == "add_numbers":
                        print("Waiting for page number position...")
                        pos_bytes = conn.recv(1024); conn.sendall(b"ACK_POSITION")
                        if not pos_bytes: raise ConnectionAbortedError("No position received.")
                        options["position"] = pos_bytes.decode()
                        print(f"Received position: {options['position']}")
                    # --- End Receiving Options ---

                    # --- Execute Action ---
                    print(f"--- Starting Action: {action} ---")
                    if action == "convert":
                        handle_file_conversion(input_ext, input_path, output_path)
                    elif action == "encrypt":
                        encrypt_pdf(input_path, output_path, options["password"])
                    elif action == "decrypt":
                        decrypt_pdf(input_path, output_path, options["password"])
                    elif action == "pdf_to_jpg":
                        created_files = convert_pdf_to_jpg(input_path, temp_dir_for_action)
                        print(f"Zipping {len(created_files)} JPGs into: {output_path}")
                        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                            for file in created_files: zipf.write(file, os.path.basename(file))
                        print(f"Finished zipping JPGs.")
                    elif action == "pdf_to_word":
                        convert_pdf_to_word(input_path, output_path)
                    elif action == "pdf_to_pptx":
                         convert_pdf_to_pptx(input_path, output_path, temp_dir_for_action)
                    elif action == "compress":
                        compress_pdf(input_path, output_path)
                    elif action == "split":
                        created_files = split_pdf(input_path, temp_dir_for_action, options["ranges"])
                        print(f"Zipping {len(created_files)} split PDFs into: {output_path}")
                        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                            for file in created_files: zipf.write(file, os.path.basename(file))
                        print(f"Finished zipping split PDFs.")
                    elif action == "merge":
                         merge_pdfs(received_zip_path, output_path)
                    elif action == "rotate":
                        rotate_pdf(input_path, output_path, options["pages"], options["angle"])
                    elif action == "add_numbers":
                        add_page_numbers_to_pdf(input_path, output_path, options["position"])
                    else:
                        raise ValueError(f"Invalid action received from client: {action}")

                    action_success = True # Mark action as successful
                    print(f"--- Action {action} Completed Successfully ---")
                    # --- End Execute Action ---

                    # --- Send Result ---
                    if not os.path.exists(output_path):
                         raise FileNotFoundError(f"Output file was not created by action '{action}': {output_path}")

                    # 5. Send Output Filename Suggestion
                    conn.sendall(output_filename_suggestion.encode())
                    ack = conn.recv(1024) # ACK_OUT_FILENAME
                    if not ack or ack != b'ACK_OUT_FILENAME': raise ConnectionAbortedError("Client disconnected before ACK filename")

                    # 6. Send Output File Size and Data
                    output_size = os.path.getsize(output_path)
                    conn.sendall(str(output_size).encode().ljust(16))
                    ack = conn.recv(1024) # ACK_OUT_SIZE
                    if not ack or ack != b'ACK_OUT_SIZE': raise ConnectionAbortedError("Client disconnected before ACK size")

                    print(f"Sending {output_size} bytes of {output_filename_suggestion} to {addr}")
                    with open(output_path, "rb") as f:
                        bytes_sent = 0
                        while True:
                            chunk = f.read(65536) # Larger chunk size for sending
                            if not chunk: break
                            conn.sendall(chunk)
                            bytes_sent += len(chunk)
                    print(f"Finished sending {bytes_sent} bytes.")
                    # --- End Send Result ---

                except (ConnectionAbortedError, ConnectionResetError, socket.timeout) as net_err:
                     print(f"! Network Error for {addr} during action '{action}': {net_err}")
                except FileNotFoundError as fnfe:
                     print(f"! File Not Found Error for {addr} during action '{action}': {fnfe}")
                     # Attempt to send specific error
                     try:
                         conn.sendall("error_file_not_found.bin".encode()); conn.recv(1024)
                         conn.sendall(str(0).encode().ljust(16)); conn.recv(1024)
                         conn.sendall(f"ERROR: File not found on server. {fnfe}".encode())
                     except Exception as e_send: print(f"Failed to send file not found error to client: {e_send}")
                except Exception as e:
                    print(f"!!! Error processing request from {addr} for action '{action}': {e}")
                    traceback.print_exc() # Print full traceback
                    if not action_success: # Only send general error if action itself failed
                        try:
                            # Send general error signal
                            conn.sendall("error_processing.bin".encode()); conn.recv(1024)
                            conn.sendall(str(0).encode().ljust(16)); conn.recv(1024)
                            error_msg_client = f"ERROR: Server failed during action '{action}'. Check server logs. Details: {str(e)[:200]}"
                            conn.sendall(error_msg_client.encode())
                        except Exception as e_send: print(f"Failed to send processing error to client: {e_send}")

                finally:
                    # --- Cleanup ---
                    print("--- Cleaning up temporary files ---")
                    paths_to_remove = [input_path, output_path]
                    dirs_to_remove = [temp_dir_for_action]
                    # Add specific merge cleanup path if action was merge
                    if action == "merge" and received_zip_path and received_zip_path != input_path:
                         paths_to_remove.append(received_zip_path)

                    for path in paths_to_remove:
                        if path and os.path.exists(path) and os.path.isfile(path):
                            try: os.remove(path); print(f"Removed file: {path}")
                            except OSError as e_rem: print(f"Warning: Error removing file {path}: {e_rem}")
                    for dir_path in dirs_to_remove:
                        if dir_path and os.path.exists(dir_path) and os.path.isdir(dir_path):
                             try: shutil.rmtree(dir_path, ignore_errors=True); print(f"Removed directory: {dir_path}")
                             except Exception as e_rem_dir: print(f"Warning: Error removing dir {dir_path}: {e_rem_dir}")
                    # --- End Cleanup ---
                    print(f"--- Connection with {addr} closed ---")

if __name__ == '__main__':
    try:
        server_program()
    except KeyboardInterrupt:
        print("\nServer stopping manually.")
    finally:
        print("Server stopped.")